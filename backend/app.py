from flask import Flask, request, jsonify, session
from flask_jwt_extended import JWTManager, jwt_required, create_access_token, decode_token, get_jwt_identity
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from flask_cors import CORS
from pptx import Presentation
import os
import quizGenrator
import requests
import random
import string
from pymongo import MongoClient
from bson import ObjectId
from datetime import datetime
from spire.doc import *
from spire.doc.common import *
from datetime import timedelta
import asyncio
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

app = Flask(__name__)
app.secret_key = "your_secret_key"
CORS(app, supports_credentials=True)

print("Connected to the server")
print("Connected to the mongodb")

app.config['JWT_ALGORITHM'] = 'HS256'
app.config["JWT_SECRET_KEY"] = "super-secret-key" 
app.config['JWT_ACCESS_TOKEN_EXPIRES'] = timedelta(minutes=30)
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['ALLOWED_EXTENSIONS'] = {'pdf','doc','docx','pptx','txt','png','jpeg','jpg'}  # Allowed file extensions

jwt = JWTManager(app)

# MongoDB client setup
client = MongoClient("mongodb+srv://harita28:harita28@cluster0.eoefb.mongodb.net/") 
db = client['Hackathon']  
quiz_collection = db['Quiz']
users_collection = db['login']

@app.route("/register", methods=["POST"])
def register():
    data = request.json
    existing_user = users_collection.find_one({"email": data["email"]})
    if existing_user:
        return jsonify({"error": "Email already exists"}), 400

    hashed_password = generate_password_hash(data["password"])
    new_user = {
        "name": data["name"],
        "email": data["email"],
        "password": hashed_password,
        "auth_provider": "email",
        "created_at": datetime.utcnow(),
    }
    users_collection.insert_one(new_user)
    return jsonify({"message": "User registered successfully!"}), 201


@app.route('/login', methods=['POST'])
def login():
    data = request.json
    email = data["email"]
    password = data["password"]

    user = users_collection.find_one({"email": email})
    if user and check_password_hash(user["password"], password):
        access_token = create_access_token(identity=email)
        return jsonify(
            access_token=access_token,
            user={
                "name": user.get("name", "User"),
                "email": user["email"]
            }
        )

    return jsonify({"message": "Invalid credentials"}), 401


@app.route('/validate-session', methods=['GET'])
@jwt_required()
def validate_session():
    current_user = get_jwt_identity()
    return jsonify({"message": f"Session valid for {current_user}"}), 200


@app.route("/get-user", methods=["GET"])
@jwt_required()
def get_user():
    current_user_email = get_jwt_identity()
    user = users_collection.find_one({"email": current_user_email})
    if user:
        return jsonify({"name": user.get("name", "User"), "email": user.get("email")})
    return jsonify({"error": "User not found"}), 404


# Generate Unique Quiz Code
def generate_quiz_code():
    return ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))

# Ensure the upload folder exists
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def allowed_file(filename):
    """Check if the file has an allowed extension."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def is_quiz_valid(quiz):
    if not isinstance(quiz, list):
        return False
    for item in quiz:
        if not isinstance(item, dict):
            return False
        if 'question' not in item or not isinstance(item['question'], str):
            return False
        if 'answer' not in item or not isinstance(item['answer'], str):
            return False
        if 'options' not in item or not isinstance(item['options'], list) or len(item['options']) != 4:
            return False
        if not all(isinstance(option, str) for option in item['options']):
            return False
    return True

@app.route('/getUserEmails', methods=['GET'])
@jwt_required()
def get_user_emails():
    try:
        users = users_collection.find({}, {"email": 1, "_id": 0})
        emails = [user['email'] for user in users]
        return jsonify({"emails": emails}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    

SENDER_EMAIL = "janhavi.22210173@viit.ac.in"
SENDER_PASSWORD = "ixxh lfpa erbu fvfa"

def send_email(recipient, subject, body):
    if not SENDER_EMAIL or not SENDER_PASSWORD:
        print("Error: Email credentials not set")
        return False
    
    try:
        msg = MIMEMultipart()
        msg['From'] = SENDER_EMAIL
        msg['To'] = recipient
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain'))

        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(SENDER_EMAIL, SENDER_PASSWORD)
            server.sendmail(SENDER_EMAIL, recipient, msg.as_string())

        print(f"Email sent to {recipient}")
        return True
    except smtplib.SMTPException as e:
        print(f"Failed to send email to {recipient}: {e}")
        return False


@app.route('/send-mails', methods=['POST'])
def send_mails():
    data = request.json
    emails = data.get('emails', [])
    custom_text = data.get('custom_text', '')

    if not emails:
        return jsonify({'error': 'No emails provided'}), 400

    subject = "Your quiz is ready!"
    predefined_body = "Hello,\n\nPlease find the attached quiz details.\n\n"

    results = {}
    for email in emails:
        body = f"{predefined_body}Custom Message: {custom_text}\n\nBest Regards,\nYour Name"
        success = send_email(email, subject, body)
        results[email] = 'Sent' if success else 'Failed'

    return jsonify(results), 200



@app.route('/upload', methods=['POST'])
@jwt_required() 
def upload_file():
    if 'files[]' not in request.files:
        return jsonify({'error': 'No files part'}), 400
    
    files = request.files.getlist('files[]')  # Get the list of files
    if not files or all(file.filename == '' for file in files):
        return jsonify({'error': 'No selected files'}), 400
    
    no_of_questions = request.form.get('noOfQuestions', None)
    question_type = request.form.get('questionType', 'mcq')
    difficulty_level = request.form.get('difficultyLevel', 'easy')
    
    combined_text = ""  # Initialize a variable to store all extracted text

    try:
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)
                print(f"File saved to {file_path}")  # Debug statement
                
                if file.filename.lower().endswith(".pdf"):
                    extracted_text = quizGenrator.extract_text_from_pdf(file_path)
                    combined_text += extracted_text + "\n"
                elif file.filename.lower().endswith((".docx" , ".doc")):
                    document = Document()   # Create a new document
                    document.LoadFromFile(file_path)   # Load the document from a file
                    extracted_text = document.GetText()   # Get the text from the document
                    combined_text += extracted_text + "\n"
                elif file.filename.lower().endswith(".txt"):
                    with open(file_path, 'r' , encoding='utf-8') as f:
                        extracted_text = f.read()
                        combined_text += extracted_text + "\n"
                elif file.filename.lower().endswith(".pptx"):
                    presentation = Presentation(file_path)
                    extracted_text = ""
                    for slide_number, slide in enumerate(presentation.slides):
                        extracted_text += f"\nSlide {slide_number + 1}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                    combined_text += extracted_text + "\n" 
                elif file.filename.lower().endswith((".png", ".jpeg", ".jpg")):
                    text = quizGenrator.get_text_from_image(file_path)
                    combined_text += text + "\n"       
                
                else:
                    print(f"Unsupported file type: {file_path}")
                
                # Extract text from the current file and append it to the combined_text variable
                 # extracted_text = quizGenrator.extract_text_from_pdf(file_path)
                #combined_text += extracted_text + "\n"  # Add a newline between texts for clarity

                # Delete the file after processing
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print(f"Deleted file: {file_path}")  # Debug statement
            else:
                return jsonify({'error': f'Invalid file type for {file.filename}, only PDFs are allowed'}), 400
        
        # Generate quiz using the combined text
        result = quizGenrator.generate_quiz_with_gemini(
            combined_text, 
            difficulty_level, 
            question_type, 
            no_of_questions
        )
        print(f"Generated quiz: {result}")  # Debug statement

        return jsonify(result)

    except Exception as e:
        print(f"Error: {e}")  # Debug statement
        return jsonify({'error': 'Error generating quiz', 'details': str(e)}), 500


@app.route('/generateGoogleForm', methods=['POST'])
@jwt_required() 
def genrateGoogleForm():
    if 'formData' not in request.json:
        return jsonify({'error': 'No form data provided'}), 400
    print("Google form ............................,",request.json);
    formData = request.json['formData']
    finalform = {'questions' : formData}
    print("Google form ............................,",finalform);

    # Validate quiz format
    if not is_quiz_valid(formData):
        return jsonify({'error': 'Invalid quiz format'}), 400

    # Process valid quiz
    apps_script_url = "https://script.google.com/macros/s/AKfycbwJEgNsRccD5G7UPs6fEphOayLZO2vYpH6NFdblrIsaXteHo1vK6UeVTh3wi5d2b7Uy/exec"
    response = requests.post(apps_script_url, json=finalform)
    if response.status_code == 200:
        print("Google Form created successfully!")
        print("Form URL:", response.json())
        return response.json()
    else:
        print("Failed to create form:", response.text)
    print(f"Result:Nothing nothing nothing")
    return jsonify(response)


@app.route('/saveQuiz', methods=['POST'])
@jwt_required() 
def save_quiz():
    data = request.json
    userId = data.get('email')
    quiz_title = data.get('quizTitle')
    quiz = data.get('quiz')
    form_url = data.get('formUrl')
    spreadsheet_url = data.get('spreadsheetUrl')
    start_date = data.get('startDate') 
    start_datetime = data.get('startDateTime')  
    end_datetime = data.get('endDateTime') 
    created_at = datetime.utcnow()
    
    if not all([quiz_title, quiz, form_url, spreadsheet_url, start_date, start_datetime, end_datetime]):
        return jsonify({"message": "Missing required fields"}), 400
    
    quiz_code = generate_quiz_code()

    if quiz_collection.find_one({"quiz_title": quiz_title}):
        return jsonify({"message": "Quiz already saved"}), 409

    quiz_collection.insert_one({
        'email': userId,
        "quiz_title": quiz_title,
        "quiz_code": quiz_code,
        "quiz": quiz,
        "formUrl": form_url,
        "spreadsheetUrl": spreadsheet_url,
        "start_date": start_date,      
        "start_datetime": start_datetime,  
        "end_datetime": end_datetime,    
        "created_at": created_at        
    })

    return jsonify({"message": "Quiz saved successfully" , 'quiz_code':quiz_code}), 201

@app.route('/join-quiz', methods=['POST'])
@jwt_required() 
def join_quiz():
    data = request.json
    quiz_code = data.get('quiz_code')

    quiz = quiz_collection.find_one({"quiz_code": quiz_code})
    if quiz:
        return jsonify({
            "quiz_title": quiz.get('quiz_title', ''),
            "questions": quiz.get('quiz', []),
            "formUrl": quiz.get('formUrl', ''),
            "spreadsheetUrl": quiz.get('spreadsheetUrl', '')
        }), 200
    return jsonify({"error": "Invalid Quiz Code"}), 404


# @app.route('/api/get-history', methods=['GET'])
# @jwt_required() 
# def get_history():

#     try:
#         current_user_email = get_jwt_identity()
#         joined_quizzes = quiz_collection.find({"email": current_user_email})
#         history = []
#         for quiz in joined_quizzes:
#             history.append({
#                 '_id': str(quiz['_id']),
#                 'quiz_title': quiz.get('quiz_title', ''),
#                 'formUrl': quiz.get('formUrl', ''),
#                 'spreadsheetUrl': quiz.get('spreadsheetUrl', ''),
#                 'joined_at': quiz.get('joined_at', datetime.now().isoformat()),
#                 'quiz_code': quiz.get('quiz_code', ''),
#                 'email': quiz.get('email', '')  # Store user's email
#             })
#         return jsonify(history), 200
#     except Exception as e:
#         print("Error fetching history:", e)
#         return jsonify({"error": "Unable to fetch history"}), 500  


@app.route('/api/get-history', methods=['GET'])
@jwt_required() 
def get_history():

    try:
        current_user_email = get_jwt_identity()
        joined_quizzes = quiz_collection.find({"email": current_user_email})
        history = []
        for quiz in joined_quizzes:
            history.append({
                '_id': str(quiz['_id']),
                'quiz_title': quiz.get('quiz_title', ''),
                'formUrl': quiz.get('formUrl', ''),
                'spreadsheetUrl': quiz.get('spreadsheetUrl', ''),
                'joined_at': quiz.get('joined_at', datetime.now().isoformat()),
                'quiz_code': quiz.get('quiz_code', ''),
                'email': quiz.get('email', ''),  # Store user's email
                'created_at': quiz.get('created_at', datetime.now().isoformat()),
            })
        return jsonify(history), 200
    except Exception as e:
        print("Error fetching history:", e)
        return jsonify({"error": "Unable to fetch history"}), 500    

from flask import jsonify, request
from flask_jwt_extended import jwt_required, get_jwt_identity
from datetime import datetime

@app.route('/api/get-joined-history', methods=['GET'])
@jwt_required()
def get_joined_history():
    try:
        current_user_email = get_jwt_identity()  # Get logged-in user's email

        # Fetch quizzes where the user's email matches the stored email in the database
        joined_quizzes = db.results.find({"email": current_user_email, "score": {"$exists": True}})

        history = []
        for quiz in joined_quizzes:
            history.append({
                '_id': str(quiz['_id']),
                'quiz_title': quiz.get('quiz_title', ''),
                'formUrl': quiz.get('formUrl', ''),
                'spreadsheetUrl': quiz.get('spreadsheetUrl', ''),
                'joined_at': quiz.get('joined_at', datetime.now().isoformat()),
                'quiz_code': quiz.get('quiz_code', ''),
                'user_email': quiz.get('user_email', ''),
                'score': quiz.get('score', '')  # Include score in response
            })

        return jsonify(history), 200
    except Exception as e:
        print("Error fetching joined quizzes:", e)
        return jsonify({"error": "Unable to fetch joined quizzes"}), 500  
    
    
@app.route('/api/get-quiz-count', methods=['GET'])
@jwt_required()
def get_quiz_count():
    try:
        # The user_id is automatically available via get_jwt_identity() since you're using @jwt_required()
        user_email = get_jwt_identity()
        print("Fetched user id:", user_email)
        if not user_email:
            return jsonify({"error": "Invalid token or missing user identity"}), 401

        # Count quizzes created by this user only
        quiz_count = quiz_collection.count_documents({'email': user_email})

        return jsonify({"quiz_count": quiz_count}), 200

    except Exception as e:
        print("Error fetching quiz count:", e)
        return jsonify({"error": "Unable to fetch quiz count"}), 500


@app.route('/api/get-quiz-details/<quiz_id>', methods=['GET'])
@jwt_required() 
def get_quiz_details(quiz_id):
    try:
        # Fetch the quiz details from MongoDB using quiz_id
        quiz = quiz_collection.find_one({"_id": ObjectId(quiz_id)})
        if quiz:
            # Return the quiz questions along with the answers
            return jsonify({
                'quiz_title': quiz.get('quiz_title', ''),
                'questions': quiz.get('quiz', [])
            }), 200
        else:
            return jsonify({'error': 'Quiz not found'}), 404
    except Exception as e:
        print("Error fetching quiz details:", e)
        return jsonify({"error": "Unable to fetch quiz details"}), 500


@app.route('/genrateQuizWithTopic', methods=['POST'])
@jwt_required() 
def genrateQuizWithTopic(): 
    # if 'topic' not in request.json:
    #     return jsonify({'error': 'No topic provided'}), 400
    no_of_questions = request.json["noOfQuestions"]
    question_type = request.json["questionType"]
    difficulty_level = request.json["difficultyLevel"]
    text = request.json['text']
    text = text.replace("\n","").replace(" ","").strip()
    # print(text)
    
    try:
        result = quizGenrator.generate_quiz_with_gemini(
            text, 
            difficulty_level, 
            question_type, 
            no_of_questions
            
        )
        print(f"Given text : {text}")
        print("question type ",question_type)
        
        print(f"Generated Quiz: {result}")
        return jsonify({'quiz': result})
    except Exception as e:
        print(f"Error generating quiz: {e}")
        return jsonify({'error': 'Failed to generate quiz'}), 500
    
@app.route('/api/quiz/<quiz_code>', methods=['GET'])
@jwt_required() 
def get_quiz(quiz_code):
    quiz = quiz_collection.find_one({'quiz_code': quiz_code})
    if quiz:
        quiz['_id'] = str(quiz['_id'])
        return jsonify(quiz)
    else:
        return jsonify({'error': 'Quiz not found'}), 404

# Submit quiz answers
@app.route('/api/submit', methods=['POST'])
@jwt_required() 
def submit_quiz():
    data = request.json
    quiz_code = data.get('quizCode')
    answers = data.get('answers')
    email = data.get('userId')

    quiz = quiz_collection.find_one({'quiz_code': quiz_code})
    if not quiz:
        return jsonify({'error': 'Quiz not found'}), 404

    score = 0
    for idx, question in enumerate(quiz['quiz']):
        correct_answer = question['answer']
        user_answer = answers.get(str(idx))
        if user_answer == correct_answer:
            score += 1

    result_data = {
        'email': email,
        'quiz_code': quiz_code,
        'score': score,
        'total': len(quiz['quiz'])
    }

    db.results.insert_one(result_data)
    return jsonify({'message': 'Quiz submitted successfully', 'score': score})

@app.route('/api/result/<quiz_code>', methods=['GET'])
@jwt_required()
def get_result(quiz_code):
    user_email = get_jwt_identity()  # assuming your identity is email or user_id from token

    result = db.results.find_one(
        {'quiz_code': quiz_code, 'email': user_email}, 
        sort=[('_id', -1)]
    )

    if result:
        result['_id'] = str(result['_id'])
        return jsonify(result)
    else:
        return jsonify({'error': 'Result not found'}), 404

    
@app.route('/add_personalized_quiz', methods=['POST'])
def add_quiz():
    try:
        data = request.json  # Receive JSON data
        quiz_code = generate_quiz_code()

        # Validate required fields
        # if not all(key in data for key in ["userId", "topic", "questions"]):
        #     return jsonify({"error": "Missing required fields: email, topic, or questions"}), 400
        if not data['questions']:
            return jsonify({"error": "Questions cannot be empty"}), 400
        if not data["userId"]:
            return jsonify({"error": "User ID cannot be empty"}), 400
        if not data["topic"]:
            return jsonify({"error": "Topic cannot be empty"}), 400

        # Create a new quiz entry with a unique ObjectId
        new_quiz = {
            "quiz_code": quiz_code,
            "email": data["userId"],
            "topic": data["topic"],
            "questions": data["questions"],
            "score": data["score"],
            "total": data["total"],
            "percentage": data["percentage"]
        }

        # Insert into MongoDB
        result = db.quiz.insert_one(new_quiz)

        return jsonify({"message": "Quiz added successfully!", "quiz_id": str(result.inserted_id)}), 201

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
@app.route('/generate-topic-name', methods=['POST'])
def generate_topic():
    try:
        data = request.json
        prompt = data.get('prompt')
        userMail = data.get('userMail')
        no_of_questions = data.get('noOfQuestions')

        if not prompt:
            return jsonify({"error": "Prompt is required"}), 400

        # Generate the topic using Gemini
        topic_name = quizGenrator.generate_name_quiz_with_gemini(prompt)
        if not topic_name:
            return jsonify({"error": "Failed to generate topic"}), 500

        # Fetch user data and extract topics
        user_data = quizGenrator.get_user_data(userMail)
        
        if not user_data:
            # No user quizzes found → Generate random quiz
            random_level = random.choice(["Easy", "Medium", "Hard"])
            print(f"No quizzes found, generating random quiz with level: {random_level}")
            response = quizGenrator.generate_personalized_quiz_with_gemini(
                prompt, random_level, question_type="mcq", no_of_questions=no_of_questions
            )
            return jsonify({
               "quiz": response,
               "topic": topic_name,
            }), 200
        
        list_of_topics = [quiz['topic'].lower() for quiz in user_data]
        topics = quizGenrator.find_similar_topic_list(topic_name, list_of_topics)

        print(f"Generated topic: {topic_name}")
        print(f"List of topics: {list_of_topics}")
        print(f"Similar topics: {topics}")

        # Filter quizzes with topics matching the generated topic
        filtered_quizzes = [
            quiz for quiz in user_data 
            if (isinstance(quiz['topic'], str) and quiz['topic'].lower() in topics) or
               (isinstance(quiz['topic'], list) and any(t.lower() in topics for t in quiz['topic']))
        ]
        
        print("Filtered quizzes:", filtered_quizzes)

        total_score = sum(quiz['percentage'] for quiz in filtered_quizzes)
        print("Total score:", total_score)
        
        if not filtered_quizzes:
            random_level = random.choice(["Easy", "Medium", "Hard"])
            print(f"No related quizzes found, generating random quiz with level: {random_level}")
            response = quizGenrator.generate_personalized_quiz_with_gemini(
                prompt, random_level, question_type="mcq", no_of_questions=no_of_questions
            )
            return jsonify({
               "quiz": response,
               "topic": topic_name,
            }), 200

        avg_total = total_score / len(filtered_quizzes)
        print("Average total:", avg_total)
        
        if avg_total >= 80:
            level = "Hard"
        elif avg_total >= 60:
            level = "Medium"
        else:
            level = "Easy"
        print("Level:", level)
            
        response = quizGenrator.generate_personalized_quiz_with_gemini(prompt, level, question_type="mcq", no_of_questions=no_of_questions)
        print(response)
        
        return jsonify({
           "quiz": response,
           "topic": topic_name,
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/getUserData/<userMail>',methods=['GET'])
def getUserData(userMail):
    try :
        results = list(db.quiz.find({"email":userMail},{"_id": 0}))
        
        if results:
            return jsonify(results)
        else:
            return jsonify({"message": "No data found for this user."}), 404
    except Exception as e :
        return jsonify({"message":"Can not get user specific data ."})
    
    
@app.route('/find_similar_topics', methods=['POST'])
def find_similar_topics():
    """
    API endpoint to find similar topics.
    Expects a JSON payload with 'topic' and 'topic_list'.
    """
    try:
        data = request.json
        topic = data.get('topics')
        user_email = data.get('userMail')
        topic_list = quizGenrator.get_user_topic_list(user_email=user_email)

        if not topic_list:
            return jsonify({'error': 'Both "topic_list" are required.'}), 400
        
        if not topic:
            return jsonify({'error': 'Both "topic" are required.'}), 400

        if not isinstance(topic_list, list):
            return jsonify({'error': '"topic_list" must be a list.'}), 400

        # Call the function to find similar topics
        similar_topics = quizGenrator.find_similar_topic_list(topic, topic_list)
        return jsonify({'similar_topics': similar_topics}), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500       
    
@app.route('/get_quiz_score', methods=['GET'])
@jwt_required()
def get_quiz_score():
    data = request.get_json()
    email = data.get('email')
    if not email:
        return jsonify({'error': 'Email is required'}), 400

    result = db.quiz.find_one
    # Fetch and return the quiz score
    return jsonify({'score': 8, 'total': 10})

    
    
if __name__ == '__main__':
    app.run(debug=True, port=5000)